"""Build a team-onboarding .pptx deck for movate-cli.

Audience: NEW users — engineers, solution architects, product leaders who
need to get from zero to "I submitted a job to production" in 15 min.

Different from `build_exec_deck.py`:
  * That deck is exec-audience: business value, status, roadmap.
  * This deck is end-user-audience: install, first agent, daily commands.

Same builder primitives (factored into both scripts independently rather
than a shared module — each deck evolves at its own pace).

Usage:
    uv run python scripts/build_onboarding_deck.py

Output: docs/movate-cli-onboarding-deck.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# -----------------------------------------------------------------------------
# Brand colors — same palette as exec deck so the two read as a set
# -----------------------------------------------------------------------------

INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_DIM = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x00, 0x6E, 0xB8)
ACCENT_LITE = RGBColor(0x00, 0xA8, 0xE8)
GREEN = RGBColor(0x2E, 0x86, 0x36)
AMBER = RGBColor(0xC1, 0x7B, 0x00)
RED = RGBColor(0xB0, 0x2A, 0x2A)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG = RGBColor(0xE6, 0xE6, 0xE6)
CODE_GREEN = RGBColor(0xA6, 0xE3, 0xA1)
CODE_BLUE = RGBColor(0x89, 0xB4, 0xFA)
BG = RGBColor(0xFF, 0xFF, 0xFF)
BG_PANEL = RGBColor(0xF6, 0xF7, 0xF9)


# -----------------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------------


def add_title_slide(prs: Presentation, title: str, subtitle: str, kicker: str) -> None:
    """Cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text(
        slide, kicker.upper(), left=0.6, top=2.2, width=12.0, height=0.5,
        size=14, bold=True, color=ACCENT, align="left",
    )
    _add_text(
        slide, title, left=0.6, top=2.7, width=12.0, height=1.6,
        size=54, bold=True, color=INK, align="left",
    )
    _add_text(
        slide, subtitle, left=0.6, top=4.3, width=12.0, height=1.5,
        size=22, color=INK_DIM, align="left",
    )
    _add_text(
        slide, "Movate internal AI agent toolkit · v1.0",
        left=0.6, top=6.6, width=12.0, height=0.5,
        size=14, color=INK_DIM, align="left",
    )


def add_section(prs: Presentation, label: str, title: str) -> None:
    """Section divider — colored band + big title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.0, top=0.0, width=13.33, height=7.5, color=BG_PANEL)
    _add_text(
        slide, label.upper(), left=0.6, top=2.8, width=12.0, height=0.5,
        size=14, bold=True, color=ACCENT, align="left",
    )
    _add_text(
        slide, title, left=0.6, top=3.3, width=12.0, height=1.5,
        size=44, bold=True, color=INK, align="left",
    )


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    subtitle: str | None = None,
    footer: str | None = None,
) -> None:
    """Generic content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(
        slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
        size=28, bold=True, color=INK, align="left",
    )
    body_top = 1.3
    if subtitle:
        _add_text(
            slide, subtitle, left=0.6, top=1.2, width=12.13, height=0.4,
            size=14, color=INK_DIM, align="left",
        )
        body_top = 1.65
    _add_bullets(
        slide, bullets, left=0.6, top=body_top, width=12.13, height=5.5,
        size=18, color=INK,
    )
    if footer:
        _add_text(
            slide, footer, left=0.6, top=6.95, width=12.13, height=0.4,
            size=11, color=INK_DIM, align="left",
        )


def add_code_slide(
    prs: Presentation,
    title: str,
    *,
    subtitle: str | None = None,
    code: str,
    explanation: list[str] | None = None,
    footer: str | None = None,
) -> None:
    """Slide with a terminal-style code block + optional bullets below.

    Use for command demonstrations — keeps the eye on the actual shell
    syntax. Code block uses a dark background for that "terminal" feel.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(
        slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
        size=28, bold=True, color=INK, align="left",
    )
    body_top = 1.3
    if subtitle:
        _add_text(
            slide, subtitle, left=0.6, top=1.2, width=12.13, height=0.4,
            size=14, color=INK_DIM, align="left",
        )
        body_top = 1.65

    # Code block — count lines to size the panel reasonably
    code_lines = code.strip().split("\n")
    code_h = max(1.4, min(3.6, 0.32 * len(code_lines) + 0.4))
    _add_filled_rect(
        slide, left=0.6, top=body_top, width=12.13, height=code_h, color=CODE_BG,
    )
    _add_code_text(
        slide, code.strip(), left=0.8, top=body_top + 0.15, width=11.7, height=code_h - 0.3,
    )

    if explanation:
        _add_bullets(
            slide, explanation,
            left=0.6, top=body_top + code_h + 0.25, width=12.13,
            height=6.9 - (body_top + code_h + 0.25),
            size=15, color=INK,
        )

    if footer:
        _add_text(
            slide, footer, left=0.6, top=6.95, width=12.13, height=0.4,
            size=11, color=INK_DIM, align="left",
        )


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    subtitle: str | None = None,
    col_widths: list[float] | None = None,
) -> None:
    """Tabular layout. Headers + N data rows."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, left=0.6, top=1.05, width=12.13, height=0.04, color=ACCENT)
    _add_text(
        slide, title, left=0.6, top=0.4, width=12.13, height=0.7,
        size=28, bold=True, color=INK, align="left",
    )
    if subtitle:
        _add_text(
            slide, subtitle, left=0.6, top=1.2, width=12.13, height=0.4,
            size=14, color=INK_DIM, align="left",
        )

    table_top = Inches(1.7 if subtitle else 1.4)
    table = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=len(headers),
        left=Inches(0.6),
        top=table_top,
        width=Inches(12.13),
        height=Inches(0.5 + 0.5 * len(rows)),
    ).table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        tf = cell.text_frame
        tf.text = header
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = BG

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_PANEL if r % 2 == 0 else BG
            tf = cell.text_frame
            tf.text = val
            p = tf.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = INK


# -----------------------------------------------------------------------------
# Primitives
# -----------------------------------------------------------------------------


def _add_text(
    slide, text: str, *, left: float, top: float, width: float, height: float,
    size: int = 16, bold: bool = False, italic: bool = False,
    color: RGBColor = INK, align: str = "left",
) -> None:
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    if align == "center":
        from pptx.enum.text import PP_ALIGN  # noqa: PLC0415
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        from pptx.enum.text import PP_ALIGN  # noqa: PLC0415
        p.alignment = PP_ALIGN.RIGHT


def _add_bullets(
    slide, items: list[str], *, left: float, top: float, width: float, height: float,
    size: int, color: RGBColor,
) -> None:
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
        p.level = 0


def _add_code_text(
    slide, code: str, *, left: float, top: float, width: float, height: float,
) -> None:
    """Monospace code text inside a code block. Uses Consolas / Menlo /
    Courier New cascade — present on every modern OS."""
    tx = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = code
    for p in tf.paragraphs:
        p.font.name = "Menlo"
        p.font.size = Pt(14)
        p.font.color.rgb = CODE_FG
        p.space_after = Pt(0)


def _add_filled_rect(
    slide, *, left: float, top: float, width: float, height: float, color: RGBColor
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


# -----------------------------------------------------------------------------
# Deck content
# -----------------------------------------------------------------------------


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ──────────────────────────────────────────────────────────────
    # Slide 1 — Title
    add_title_slide(
        prs,
        kicker="Team onboarding",
        title="movate-cli",
        subtitle="From zero to first deployed AI agent in 15 minutes",
    )

    # Slide 2 — What is movate-cli
    add_content_slide(
        prs,
        title="What is movate-cli?",
        subtitle="The one-paragraph answer",
        bullets=[
            "A command-line toolkit for building, evaluating, and deploying AI agents — the same way you'd build a typed Python library.",
            "Declarative `agent.yaml` (typed JSON in, typed JSON out). Prompts are version-controlled markdown. Evals are dataset-based with regression gates.",
            "Multi-vendor LLM support — OpenAI, Anthropic, Gemini, Azure OpenAI — behind one Protocol. Models are config; swap freely without rewriting code.",
            "One CLI binary covers the full loop: scaffold → iterate → eval → deploy → submit jobs. Same tool on your laptop and in CI.",
            "Movate's deployed runtime on Azure handles production traffic. You submit jobs from your terminal (or your code) and get notified when they're done.",
        ],
    )

    # Slide 3 — What you can do with it
    add_content_slide(
        prs,
        title="What you can build",
        subtitle="Concrete capabilities",
        bullets=[
            "**Typed AI agents** — `agent.yaml` declares the input + output schema; movate enforces them on every call. No more \"the model returned a different shape today.\"",
            "**Eval suites with CI gates** — `movate eval --gate 0.7` fails CI when quality drops. `--baseline <id>` detects drift vs prior runs.",
            "**Cost + latency comparisons** — `movate bench` runs the same prompt across N providers; reports cost mean, latency p50/p95, score under your gate mode.",
            "**Multi-step workflows** — chain agents with conditional branches, parallel fan-out, human-in-the-loop pauses. Optional LangGraph backend for advanced topologies.",
            "**Production deployment** — `movate deploy` builds + ships to Movate's Azure Container Apps. Per-tenant cost ceilings, rate limiting, tenant isolation all built-in.",
            "**Personal notifications** — Telegram bot pings your phone when your submitted jobs land terminal. Free, 5-minute setup.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # Slide 4 — Section: INSTALL
    add_section(prs, "Part 1", "Install (3 commands, ~2 minutes)")

    # Slide 5 — Prereqs
    add_content_slide(
        prs,
        title="Before you start",
        subtitle="What you need",
        bullets=[
            "**macOS, Linux, or Windows (WSL).** No specific OS requirement.",
            "**A terminal.** Anything modern: iTerm2, Terminal.app, Windows Terminal, VS Code's integrated terminal.",
            "**No Python install needed.** We use uv, which manages Python for you.",
            "**No Docker, no virtualenvs to manage yourself, no PATH tweaking.** uv handles all of it.",
            "**~5 minutes of attention.** Total install + first agent run.",
        ],
        footer="If you have a corporate proxy / VPN, check it can reach github.com and astral.sh — uv installs from both.",
    )

    # Slide 6 — Install command 1: uv
    add_code_slide(
        prs,
        title="Step 1 — Install uv",
        subtitle="One-line installer; takes ~10 seconds",
        code="""# macOS / Linux / WSL
curl -LsSf https://astral.sh/uv/install.sh | sh

# Windows (PowerShell)
powershell -ExecutionPolicy ByPass -c "irm https://astral.sh/uv/install.ps1 | iex"

# Verify
uv --version""",
        explanation=[
            "uv is a Python package manager — think `pip` + `pyenv` + `pipx` fused into one fast binary.",
            "Runs on every modern OS. Installs to `~/.local/bin/uv`; the installer adds that to your PATH.",
            "Open a new terminal after install so PATH refresh takes effect, OR run `source ~/.zshrc` (or your shell rc).",
        ],
    )

    # Slide 7 — Install command 2: movate-cli
    add_code_slide(
        prs,
        title="Step 2 — Install movate-cli",
        subtitle="One command; pulls the latest pinned release from GitHub",
        code="""# Install from the Movate GitHub release
uv tool install --from git+https://github.com/Movate/movate-cli@v1.0.0 movate-cli

# Verify it's on your PATH
movate --version
movate --help""",
        explanation=[
            "`uv tool install` puts movate in an isolated venv but exposes the `movate` binary on your PATH globally.",
            "Pinned to `@v1.0.0` for stability. To upgrade later: `uv tool upgrade movate-cli`.",
            "`movate --help` shows commands grouped into panels (Develop / Run & evaluate / Diagnose / Deploy & operate / Manage) — that's the full command surface.",
        ],
        footer="If `movate` isn't found after install, run `uv tool update-shell` and restart your terminal.",
    )

    # Slide 8 — Connect to Movate's runtime
    add_code_slide(
        prs,
        title="Step 3 — Connect to Movate's runtime",
        subtitle="Get your API key + register the deployment target",
        code="""# Grab your API key from 1Password
#   → entry: "movate-toolkit-prod"
#   → copy the value (looks like: mvt_live_acme_abc12345_<long-secret>)

# Paste into your shell (silent: no echo, no shell history)
read -rs MOVATE_PROD_KEY && export MOVATE_PROD_KEY

# Register the deployment target
movate config add-target prod \\
    --url https://movate-prod-api.<region>.azurecontainerapps.io \\
    --key-env MOVATE_PROD_KEY \\
    --set-active

# Smoke test: 9 checks, all should be green
movate doctor --target prod""",
        explanation=[
            "Your key lives in the env var `MOVATE_PROD_KEY` — never in `~/.movate/config.yaml`. Each shell session needs it re-exported (consider adding to your `.zshrc`).",
            "`movate doctor` walks: az login → subscription → resource group → ACR → both Container Apps → `/healthz`. If anything's red, the output tells you the fix.",
        ],
        footer="Replace <region> with the prod runtime's actual URL — your team lead has it; also in 1Password.",
    )

    # ──────────────────────────────────────────────────────────────
    # Slide 9 — Section: BUILD YOUR FIRST AGENT
    add_section(prs, "Part 2", "Build your first agent")

    # Slide 10 — Scaffold
    add_code_slide(
        prs,
        title="Day 0: Scaffold a new agent",
        subtitle="`movate init` creates the typed skeleton + sample dataset",
        code="""movate init my-faq-agent -t faq
cd my-faq-agent
tree .""",
        explanation=[
            "`-t faq` picks the FAQ-bot template; alternatives: `summarizer`, `classifier`, default `agent_init`.",
            "Creates: `agent.yaml` (the contract), `prompt.md` (system prompt), `schema/{input,output}.json` (typed I/O), `evals/dataset.jsonl` (test cases), `evals/judge.yaml.example` (quality scoring config).",
            "Every file is plain text + version-controllable. No magic; no opaque binaries.",
        ],
    )

    # Slide 11 — agent.yaml anatomy
    add_code_slide(
        prs,
        title="The contract: `agent.yaml`",
        subtitle="One declarative file describes the whole agent",
        code="""api_version: movate/v1
kind: Agent
name: my-faq-agent
version: 0.1.0
description: Answers FAQs from a small corpus.

model:
  provider: openai/gpt-4o-mini-2024-07-18
  params:
    temperature: 0.2
  fallback:
    - provider: anthropic/claude-haiku-4-5-20251001

prompt: ./prompt.md
schema:
  input: ./schema/input.json
  output: ./schema/output.json

budget:
  max_cost_usd_per_run: 0.05""",
        explanation=[
            "`api_version` pins the YAML schema — bumping it is a deliberate migration event, never accidental.",
            "`model.fallback` is a chain — if the primary errors out (rate limit, transient failure), movate tries each fallback in order.",
            "`budget.max_cost_usd_per_run` is a hard cap. A run that would exceed it aborts BEFORE the provider call. Zero surprise bills.",
        ],
    )

    # Slide 12 — Hot-reload TDD
    add_code_slide(
        prs,
        title="Day 0: Hot-reload while you edit",
        subtitle="`movate watch` re-validates on every save — TDD-style feedback in <1s",
        code="""# Terminal 1: hot-reload watcher
movate watch ./my-faq-agent

# Terminal 2: open prompt.md in your editor
$EDITOR ./my-faq-agent/prompt.md

# Save the file → Terminal 1 instantly shows:
#  ✓ schema valid
#  ⚠ prompt linter warnings (if any)
#  ✓ cost forecast: ~$0.042 for 30-case eval""",
        explanation=[
            "Prompt linter has 4 rules: undeclared `{{ input.X }}` refs, empty prompt, missing JSON instruction, missing output schema field references.",
            "Cost forecast estimates eval-run cost from dataset size × pricing table BEFORE you spend money running it.",
            "200ms debounce — won't fire on editor-temporary writes (vim's swap files, VS Code's incremental saves).",
        ],
    )

    # Slide 13 — Run locally
    add_code_slide(
        prs,
        title="Day 0: Test a single run",
        subtitle="`movate run` — real LLM call or deterministic mock",
        code="""# Real LLM call (uses OPENAI_API_KEY from .env)
movate run ./my-faq-agent "What is movate?"

# Mock provider — deterministic, no API key, no $ — for smoke tests
movate run ./my-faq-agent "test input" --mock

# Output formats: json (default, pipe-friendly) | text (human)
movate run ./my-faq-agent "test" --mock --output text""",
        explanation=[
            "JSON output is pipe-friendly: `movate run ... | jq '.output.answer'`.",
            "Mock mode is for CI / smoke testing — no API calls, deterministic responses, zero cost.",
            "Every run is persisted locally to `~/.movate/local.db` (sqlite) — replay later with `movate trace replay <run-id>`.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # Slide 14 — Section: QUALITY GATES
    add_section(prs, "Part 3", "Quality gates (eval + bench)")

    # Slide 15 — Eval
    add_code_slide(
        prs,
        title="Day 1: Eval your agent",
        subtitle="`movate eval` — regression-gated dataset evaluation",
        code="""# Eval against evals/dataset.jsonl with the configured judge
movate eval ./my-faq-agent --gate 0.7 --runs 3

# Output: Rich table of per-case results + aggregate score
# Exit 1 if mean score < 0.7 — CI-gateable

# Drift detection: diff this run vs a stored baseline
movate eval ./my-faq-agent --baseline <eval-id> \\
    --regression-tolerance 0.05""",
        explanation=[
            "Exact-match scorer for deterministic test cases (e.g., classifiers); LLM-as-judge for free-form responses.",
            "Cross-family enforcement: judge model can't share vendor family with tested model (catches confounded scores).",
            "`--runs 3` smooths variance across runs; `--gate-mode mean | min | p10` controls how N runs aggregate.",
            "Every eval is persisted as an `EvalRecord` — pass `--baseline <eval-id>` later to detect drift; CI exits 1 on regression past tolerance.",
        ],
    )

    # Slide 16 — Bench
    add_code_slide(
        prs,
        title="Day 1: Compare models",
        subtitle="`movate bench` — same input across N providers, with cost/latency/quality",
        code="""# Compare 3 models on the same input
movate bench ./my-faq-agent "what is movate?" --runs 3 \\
    -m openai/gpt-4o-mini-2024-07-18 \\
    -m anthropic/claude-haiku-4-5-20251001 \\
    -m google/gemini-2-flash

# Per-model report: cost mean, latency p50/p95, score (if judge configured)
# Saved as a BenchRecord with bench_id for later --baseline diffs""",
        explanation=[
            "Each provider gets N runs (variance smoothing); reported as cost mean, latency p50/p95, score under the configured gate mode.",
            "Cross-family judge enforcement: a judge model from family X gets skipped against any tested model in family X.",
            "`--baseline <bench-id>` shows per-model deltas vs a stored bench — score, cost, latency, model-set drift. CI-gateable on score regression past `--regression-tolerance`.",
        ],
    )

    # ──────────────────────────────────────────────────────────────
    # Slide 17 — Section: USE THE DEPLOYED RUNTIME
    add_section(prs, "Part 4", "Use the deployed runtime")

    # Slide 18 — Submit + jobs
    add_code_slide(
        prs,
        title="Day 2+: Submit a job",
        subtitle="`movate submit` queues against the deployed runtime",
        code="""# Fire-and-forget — returns job_id immediately
movate submit my-faq-agent '{"question": "what is movate?"}' --target prod

# --wait blocks until terminal; prints final status
movate submit my-faq-agent '{"question": "..."}' --target prod --wait

# Check a specific job
movate jobs show <job-id>

# Recent jobs (filter by status / agent / target)
movate jobs list --target prod --status success --limit 10""",
        explanation=[
            "`--target prod` selects which deployed runtime (you configured earlier). Per-target bearer token from env var.",
            "Jobs queue server-side; the worker claims them, runs the agent, persists the result. You poll via `jobs show` or wait via `--wait`.",
            "Multi-tenant safe — your `tenant_id` is encoded in your API key. You can never see another tenant's jobs.",
        ],
    )

    # Slide 19 — Telegram
    add_content_slide(
        prs,
        title="Day 2+: Get notified when jobs land",
        subtitle="Optional — 5-minute Telegram bot setup (worth it)",
        bullets=[
            "**Why Telegram?** Free, instant, works on phone/desktop/web. No SMS regulatory hassle. Notification looks like: *\"✅ movate agent/my-faq-agent — success (423ms) — run: 5fdb30de\"*.",
            "**Setup (5 min, one-time):**",
            "  1. In Telegram, message `@BotFather` → send `/newbot` → save the token.",
            "  2. Search for your new bot in Telegram → tap `Start`.",
            "  3. Open `https://api.telegram.org/bot<TOKEN>/getUpdates` in browser → copy your `chat.id`.",
            "  4. Ask your team lead to add your token to Key Vault + your chat_id to the deployed runtime config (one Bicep redeploy).",
            "  5. Done. Every job you submit pings your phone within seconds of completion.",
            "**Detailed runbook**: `docs/azure-bootstrap.md` → section *\"Optional: Telegram alerts\"*.",
        ],
        footer="Email + SMS channels also available (`--notify-email`, `--notify-sms`) for per-job opt-in; Telegram is operator-wide.",
    )

    # ──────────────────────────────────────────────────────────────
    # Slide 20 — Section: REFERENCE
    add_section(prs, "Part 5", "Reference: cheat sheet + troubleshooting")

    # Slide 21 — Cheat sheet
    add_table_slide(
        prs,
        title="Common tasks cheat sheet",
        subtitle="When you want to... run this",
        headers=["I want to...", "Command"],
        rows=[
            ["Start a new agent", "`movate init <name> -t <template>`"],
            ["See what `-t` templates exist", "`movate init --list-templates`"],
            ["Iterate fast on a prompt", "`movate watch ./<agent>`"],
            ["Run once (real LLM)", "`movate run ./<agent> \"input\"`"],
            ["Run once (mock, no $)", "`movate run ./<agent> \"input\" --mock`"],
            ["Eval against a dataset", "`movate eval ./<agent> --gate 0.7 --runs 3`"],
            ["Detect eval regression", "`movate eval ./<agent> --baseline <eval-id>`"],
            ["Compare 3 models", "`movate bench ./<agent> \"input\" -m m1 -m m2 -m m3`"],
            ["Submit to production", "`movate submit <agent> '{...}' --target prod`"],
            ["Check a job", "`movate jobs show <job-id>`"],
            ["List recent jobs", "`movate jobs list --status success --limit 10`"],
            ["Replay a past run", "`movate run ./<agent> --replay <run-id>`"],
            ["Health-check the deployment", "`movate doctor --target prod`"],
            ["Print pricing table", "`movate pricing`"],
            ["Get help on any command", "`movate <command> --help`"],
        ],
        col_widths=[5.5, 6.63],
    )

    # Slide 22 — Troubleshooting
    add_table_slide(
        prs,
        title="Troubleshooting",
        subtitle="Common errors → the fix",
        headers=["Symptom", "Likely cause", "Fix"],
        rows=[
            [
                "`movate: command not found`",
                "PATH not refreshed after install",
                "Restart terminal, or run `uv tool update-shell` then restart",
            ],
            [
                "`movate doctor` shows red on subscription",
                "Logged in to wrong Azure sub",
                "`az account set --subscription <id>`",
            ],
            [
                "`401 Unauthorized` on `movate submit`",
                "API key not in env, or wrong key",
                "`echo \"len=${#MOVATE_PROD_KEY}\"` — expect ~70. Re-export if 0.",
            ],
            [
                "`movate eval --gate 0.7` fails on what should pass",
                "Judge is in same family as tested model",
                "Pick a judge from a different vendor in `evals/judge.yaml`",
            ],
            [
                "Schema validation error on input",
                "Input dict doesn't match `schema/input.json`",
                "Open the schema file, check required fields + types",
            ],
            [
                "`movate watch` keeps re-firing on no change",
                "Editor doing atomic-save (write to temp + rename)",
                "Set `--poll-interval 1.0` to debounce more aggressively",
            ],
            [
                "Cost forecast says $40 for an eval",
                "Dataset too large OR temperature too high",
                "`head -10 evals/dataset.jsonl > evals/dataset-small.jsonl`; iterate first",
            ],
            [
                "Job submitted but no Telegram ping",
                "Either bot not in env, or chat is muted in Telegram",
                "Worker logs: `az containerapp logs show ... | grep telegram`",
            ],
        ],
        col_widths=[3.5, 4.0, 4.63],
    )

    # Slide 23 — Where to get help
    add_content_slide(
        prs,
        title="Where to get help",
        subtitle="In order of fastest-resolution",
        bullets=[
            "**`movate <command> --help`** — every command has detailed help with examples. Often answers your question without leaving the terminal.",
            "**`movate doctor --target <env>`** — 9-check diagnostic against your deployed runtime. Tells you what's broken AND how to fix it.",
            "**Repository docs** — `docs/dev-loop.md` (this content, deeper), `docs/azure-bootstrap.md` (Azure-specific setup), `docs/v1.0-overview.md` (architecture).",
            "**Internal Slack: #movate-cli** — fastest for \"is this the right approach?\" questions or quick unblocks.",
            "**GitHub Issues**: [github.com/Movate/movate-cli/issues](https://github.com/Movate/movate-cli/issues) — for bugs, feature requests, or anything where the conversation needs a thread.",
            "**Office hours** — Wednesdays 2-3pm PT, drop in any time. (Or ping your team lead.)",
        ],
        footer="When reporting issues: include `movate --version`, your OS, the exact command, and the full error output.",
    )

    # Slide 24 — What's possible next
    add_content_slide(
        prs,
        title="What to explore next",
        subtitle="Once you're comfortable with the basics",
        bullets=[
            "**Multi-step workflows** — chain agents with `workflow.yaml`. Linear DAGs out of the box; conditional / parallel / human-in-the-loop via `runtime: langgraph`.",
            "**Tool-calling agents** — register Python functions as tools the agent can invoke. `@tool` decorator handles the JSON-schema plumbing.",
            "**CI integration** — drop `.github/workflows/eval-gate.example.yml` (in the repo) into your project; eval regressions block PRs automatically.",
            "**Custom judges** — write your own `evals/judge.yaml` for domain-specific scoring rubrics.",
            "**Replay debugging** — `movate trace replay <run-id>` reconstructs a past run's call tree from stored RunRecords. Great for debugging \"this used to work\" issues.",
            "**Per-tenant budgets** — `movate tenants set-budget --monthly-usd-limit 50` caps your team's spend.",
        ],
        footer="Read `docs/dev-loop.md` for the full developer-facing flow + `docs/v1.0-overview.md` for the system architecture.",
    )

    # Slide 25 — Closing
    add_content_slide(
        prs,
        title="You're ready",
        subtitle="Recap of the 5 commands you'll run most",
        bullets=[
            "**`movate init <name> -t <template>`** — start a new agent from a known-good template.",
            "**`movate watch ./<agent>`** — iterate on the prompt with sub-second feedback.",
            "**`movate eval ./<agent> --gate 0.7`** — gate quality before merging.",
            "**`movate submit <agent> '...' --target prod`** — run it against Movate's deployed infrastructure.",
            "**`movate jobs show <job-id>`** — check what happened.",
            "Everything else is variations. `movate --help` is the discovery surface — start there when you're unsure.",
            "Good luck. Ping #movate-cli when you ship your first agent — we'd love to see it.",
        ],
    )

    return prs


if __name__ == "__main__":
    prs = build()
    out = Path("docs/movate-cli-onboarding-deck.pptx")
    prs.save(out)
    print(f"✓ wrote {out} ({len(prs.slides)} slides)")
